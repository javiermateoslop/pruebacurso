import os
from pathlib import Path
from pptx import Presentation

# Directorios
RESOURCES_DIR = Path(r"u:/antigravityprojects/libroprueba/pruebacurso/Recursos")
OUTPUT_DIR = Path(r"u:/antigravityprojects/libroprueba/pruebacurso/book/_static/recursos")

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

for pptx_file in RESOURCES_DIR.glob("*.pptx"):
    prs = Presentation(pptx_file)
    slide_idx = 1
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture
                image = shape.image
                image_bytes = image.blob
                ext = image.ext
                safe_stem = pptx_file.stem.replace(' ', '_').replace('-', '_')
                while '__' in safe_stem: safe_stem = safe_stem.replace('__', '_')
                image_name = f"{safe_stem}_slide{slide_idx}.{ext}"
                image_path = OUTPUT_DIR / image_name
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
        slide_idx += 1
print("Extracción completada.")
