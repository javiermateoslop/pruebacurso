import os
from pathlib import Path
from pptx import Presentation
from PIL import Image
import io

def extract_images_from_pptx(pptx_path: Path, output_dir: Path):
    """Extract all embedded images from a PPTX file.
    Each image is saved with its original format (png, jpg, etc.).
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(pptx_path)
    count = 0
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.shape_type == 13:  # 13 = Picture
                continue
            image = shape.image
            image_bytes = image.blob
            ext = image.ext  # e.g., 'png', 'jpeg'
            filename = f"slide_{slide_idx:03d}_{count}.{ext}"
            out_path = output_dir / filename
            with open(out_path, 'wb') as f:
                f.write(image_bytes)
            count += 1
    if count == 0:
        # No embedded images; create a placeholder PNG indicating empty slide
        placeholder = Image.new('RGB', (800, 600), color='lightgray')
        placeholder_path = output_dir / f"slide_{slide_idx:03d}_placeholder.png"
        placeholder.save(placeholder_path)

def main():
    base_dir = Path(__file__).parents[2]
    recursos_dir = base_dir / "Recursos"
    output_base = base_dir / "book" / "_static" / "recursos"
    for pptx_file in recursos_dir.glob("*.pptx"):
        out_dir = output_base / pptx_file.stem.replace(' ', '_')
        print(f"Extracting images from {pptx_file.name} -> {out_dir}")
        extract_images_from_pptx(pptx_file, out_dir)

if __name__ == "__main__":
    main()
